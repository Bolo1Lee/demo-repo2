import fitz  # PyMuPDF
from pptx import Presentation

def extraer_texto_pdf(archivo):
    """
    Extrae el texto desde un archivo subido (UploadedFile de Django).
    """
    texto = ""
    with fitz.open(stream=archivo.read(), filetype="pdf") as doc:
        for pagina in doc:
            texto += pagina.get_text()
    return texto


from pptx import Presentation

def extraer_texto_ppt(archivo):
    """
    Extrae texto desde un archivo PPT o PPTX subido por el usuario.
    """
    texto = ""
    prs = Presentation(archivo)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto += shape.text + "\n"
    return texto

